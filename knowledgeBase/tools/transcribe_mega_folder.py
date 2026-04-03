from __future__ import annotations

import argparse
import csv
import html
import json
import os
import re
import shutil
import subprocess
import sys
import threading
import time
from collections import deque
from concurrent.futures import FIRST_COMPLETED, ThreadPoolExecutor, wait
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path, PurePosixPath
from shutil import which
from tempfile import TemporaryDirectory
from typing import Iterable


SUPPORTED_MEDIA_SUFFIXES = {
    ".aac",
    ".aiff",
    ".avi",
    ".flac",
    ".m4a",
    ".m4v",
    ".mkv",
    ".mov",
    ".mp3",
    ".mp4",
    ".mpeg",
    ".mpg",
    ".ogg",
    ".wav",
    ".webm",
}
SUPPORTED_IMAGE_SUFFIXES = {
    ".bmp",
    ".gif",
    ".jpeg",
    ".jpg",
    ".png",
    ".svg",
    ".tif",
    ".tiff",
    ".webp",
}
SUPPORTED_NATIVE_DOCUMENT_SUFFIXES = {
    ".csv",
    ".docx",
    ".htm",
    ".html",
    ".json",
    ".md",
    ".pptx",
    ".rtf",
    ".srt",
    ".tsv",
    ".txt",
    ".xls",
    ".xlsx",
}
SUPPORTED_OCR_DOCUMENT_SUFFIXES = {
    ".doc",
    ".docx",
    ".pdf",
    ".ppt",
    ".pptx",
    ".xls",
    ".xlsx",
}
SUPPORTED_DOCUMENT_SUFFIXES = SUPPORTED_NATIVE_DOCUMENT_SUFFIXES | SUPPORTED_OCR_DOCUMENT_SUFFIXES
STRUCTURED_CONTENT_MODE = "structured"
TIMESTAMP_CONTENT_MODE = "timestamps"
NATIVE_PROCESSOR_PREFIX = "native"
MARKER_PROCESSOR = "marker"
MEDIA_PROCESSOR = "media_transcriber"
GENERATED_ARTIFACT_SUFFIX = ".extracted.txt"
PROCESSABLE_FILE_KINDS = ("document", "image", "media")

MEGA_DONE_SUFFIX = "_transcript_done"
MEGA_HANDLE_PREFIX = "H:"
MEGA_LONG_ROW_RE = re.compile(
    r"^(?P<flags>\S+)\s+"
    r"(?P<vers>\S+)\s+"
    r"(?P<size>\S+)\s+"
    r"(?P<date>\S+(?:\s+\S+)?)\s+"
    r"(?P<handle>H:[^\s]+)\s+"
    r"(?P<name>.+)$"
)
MEGA_TIME_FORMAT = "ISO6081_WITH_TIME"
ISO_UTC_FORMAT = "%Y-%m-%dT%H:%M:%SZ"
SHORT_UTC_FORMAT = "%d%b%Y %H:%M:%S"
MEGA_DOWNLOAD_RETRY_ATTEMPTS = 3
MEGA_DOWNLOAD_RETRY_BACKOFF_SECONDS = 2
MAX_TRANSCRIBE_WORKERS = 32


@dataclass(frozen=True)
class SourceSpec:
    source_kind: str
    source_value: str
    output_root: PurePosixPath
    source_label: str
    source_root: PurePosixPath | None = None


@dataclass(frozen=True)
class BatchItem:
    source_identifier: str
    relative_path: PurePosixPath
    source_file_kind: str
    output_relative_path: PurePosixPath
    companion_output_relative_path: PurePosixPath
    output_path: PurePosixPath
    companion_output_path: PurePosixPath
    mega_source_path: PurePosixPath | None = None
    local_source_path: Path | None = None


@dataclass(frozen=True)
class MegaLsEntry:
    path: PurePosixPath
    parent_path: PurePosixPath
    name: str
    kind: str
    handle: str | None = None
    created_at_utc: str | None = None
    modified_at_utc: str | None = None


@dataclass(frozen=True)
class BatchRuntimeItem:
    batch_item: BatchItem
    local_source_path: Path
    local_output_path: Path
    local_companion_output_path: Path
    staging_key: str | None
    output_key: str | None
    companion_output_key: str | None
    manifest_index: int


@dataclass(frozen=True)
class ProcessResult:
    processor: str
    content_mode: str
    metadata: dict[str, str] | None = None


@dataclass(frozen=True)
class ResumeDecision:
    status: str
    reason: str
    queue_for_processing: bool


def utc_timestamp() -> str:
    return datetime.utcnow().strftime("%Y%m%d-%H%M%S")


def normalize_mega_path(value: str | PurePosixPath) -> PurePosixPath:
    path = PurePosixPath(str(value).strip())
    if not str(path).startswith("/"):
        raise ValueError("MEGA paths must be absolute paths starting with '/'.")
    return path


def default_output_root(source_root: str | PurePosixPath) -> PurePosixPath:
    return normalize_mega_path(source_root)


def classify_source_file_kind(path: str | PurePosixPath) -> str:
    candidate = PurePosixPath(path)
    suffix = candidate.suffix.lower()
    if suffix in SUPPORTED_MEDIA_SUFFIXES:
        return "media"
    if suffix in SUPPORTED_IMAGE_SUFFIXES:
        return "image"
    if suffix in SUPPORTED_DOCUMENT_SUFFIXES:
        return "document"
    return "other"


def normalize_selected_file_kinds(
    include_file_kinds: str | Iterable[str] | None,
) -> set[str]:
    if include_file_kinds is None:
        return set(PROCESSABLE_FILE_KINDS)
    if isinstance(include_file_kinds, str):
        raw_values = include_file_kinds.split(",")
    else:
        raw_values = list(include_file_kinds)
    selected = {str(value).strip() for value in raw_values if str(value).strip()}
    invalid = selected.difference(PROCESSABLE_FILE_KINDS)
    if invalid:
        raise ValueError(
            "Unsupported file kind selection: "
            + ", ".join(sorted(invalid))
            + ". Expected a subset of: "
            + ", ".join(PROCESSABLE_FILE_KINDS)
        )
    if not selected:
        raise ValueError("At least one file kind must be selected.")
    return selected


def normalize_preferred_regions(
    preferred_regions: str | Iterable[str] | None,
) -> list[str]:
    if preferred_regions is None:
        return []
    if isinstance(preferred_regions, str):
        raw_values = preferred_regions.split(",")
    else:
        raw_values = list(preferred_regions)
    return [str(value).strip() for value in raw_values if str(value).strip()]


def is_generated_companion_path(path: str | PurePosixPath) -> bool:
    name = PurePosixPath(path).name.lower()
    return name.endswith(".txt") and (".structured." in name or ".timestamps." in name)


def is_generated_pipeline_artifact(path: str | PurePosixPath) -> bool:
    candidate = PurePosixPath(path)
    name = candidate.name.lower()
    return (
        name.endswith(GENERATED_ARTIFACT_SUFFIX)
        or is_generated_companion_path(candidate)
    )


def mega_path_is_within(path: str | PurePosixPath, root: str | PurePosixPath) -> bool:
    normalized_path = normalize_mega_path(path)
    normalized_root = normalize_mega_path(root)
    return normalized_path == normalized_root or str(normalized_path).startswith(f"{normalized_root}/")


def build_companion_prefix(relative_path: PurePosixPath, *, source_file_kind: str) -> str:
    marker = ".timestamps." if source_file_kind == "media" else ".structured."
    return f"{relative_path.stem}{marker}"


def is_expected_companion_output(
    candidate: str | PurePosixPath,
    batch_item: BatchItem,
) -> bool:
    candidate_path = normalize_mega_path(candidate)
    return (
        candidate_path.parent == batch_item.companion_output_path.parent
        and candidate_path.name.startswith(
            build_companion_prefix(
                batch_item.relative_path,
                source_file_kind=batch_item.source_file_kind,
            )
        )
        and candidate_path.suffix.lower() == ".txt"
    )


def find_existing_companion_outputs(
    batch_item: BatchItem,
    remote_paths: set[PurePosixPath],
) -> list[PurePosixPath]:
    matches = [
        candidate
        for candidate in remote_paths
        if is_expected_companion_output(candidate, batch_item)
    ]
    return sorted(matches)


def find_stale_generated_outputs(
    batch_item: BatchItem,
    remote_paths: set[PurePosixPath],
) -> list[PurePosixPath]:
    prefix = f"{batch_item.relative_path.stem}."
    stale_paths: list[PurePosixPath] = []
    for candidate in sorted(remote_paths):
        if candidate.parent != batch_item.output_path.parent:
            continue
        if candidate == batch_item.output_path:
            continue
        if is_expected_companion_output(candidate, batch_item):
            continue
        if not candidate.name.startswith(prefix):
            continue
        if is_generated_pipeline_artifact(candidate):
            stale_paths.append(candidate)
    return stale_paths


def classify_resume_decision(
    batch_item: BatchItem,
    remote_paths: set[PurePosixPath],
    *,
    force: bool,
) -> ResumeDecision:
    if force:
        return ResumeDecision(
            status="pending",
            reason="force rerun requested; live MEGA outputs were not used for skipping",
            queue_for_processing=True,
        )

    has_canonical_output = batch_item.output_path in remote_paths
    existing_companion_outputs = find_existing_companion_outputs(batch_item, remote_paths)
    has_companion_output = bool(existing_companion_outputs)
    stale_generated_outputs = find_stale_generated_outputs(batch_item, remote_paths)

    if has_canonical_output and has_companion_output:
        return ResumeDecision(
            status="complete",
            reason="found canonical and companion outputs on MEGA",
            queue_for_processing=False,
        )
    if has_canonical_output:
        return ResumeDecision(
            status="partial_missing_companion",
            reason="found canonical output on MEGA but companion output is missing",
            queue_for_processing=True,
        )
    if has_companion_output:
        return ResumeDecision(
            status="partial_missing_canonical",
            reason="found companion output on MEGA but canonical output is missing",
            queue_for_processing=True,
        )
    if stale_generated_outputs:
        return ResumeDecision(
            status="stale_generated_outputs",
            reason="found stale generated artifacts on MEGA without a complete canonical and companion pair",
            queue_for_processing=True,
        )
    return ResumeDecision(
        status="pending",
        reason="no existing generated outputs found on MEGA",
        queue_for_processing=True,
    )


def update_resume_summary(
    summary: dict[str, int],
    decision: ResumeDecision,
) -> None:
    if not decision.queue_for_processing and decision.status == "complete":
        summary["completed_skipped"] += 1
        return
    if decision.status == "retry_failed":
        summary["failed_requeued"] += 1
        return
    if decision.status in {
        "partial_missing_companion",
        "partial_missing_canonical",
        "stale_generated_outputs",
    }:
        summary["partial_requeued"] += 1
        return
    summary["pending_queued"] += 1


def build_remote_resume_paths(
    remote_paths: Iterable[str | PurePosixPath],
    *,
    root: str | PurePosixPath,
) -> set[PurePosixPath]:
    normalized_root = normalize_mega_path(root)
    return {
        normalize_mega_path(path)
        for path in remote_paths
        if mega_path_is_within(path, normalized_root)
    }


def media_transcript_source_for_path(
    candidate: PurePosixPath,
    all_paths: set[PurePosixPath],
) -> PurePosixPath | None:
    if candidate.suffix.lower() != ".txt":
        return None
    for suffix in SUPPORTED_MEDIA_SUFFIXES:
        source_candidate = candidate.with_suffix(suffix)
        if source_candidate in all_paths:
            return source_candidate
    return None


def append_name_before_suffix(path: PurePosixPath, inserted_suffix: str) -> PurePosixPath:
    return path.with_name(f"{path.stem}{inserted_suffix}")


def build_output_path(
    source_file: str | PurePosixPath,
    source_root: str | PurePosixPath,
    output_root: str | PurePosixPath | None = None,
) -> PurePosixPath:
    return build_generated_output_path(
        source_file,
        source_root,
        source_file_kind="media",
        output_root=output_root,
    )


def build_generated_output_path(
    source_file: str | PurePosixPath,
    source_root: str | PurePosixPath,
    *,
    source_file_kind: str,
    output_root: str | PurePosixPath | None = None,
) -> PurePosixPath:
    source_root_path = normalize_mega_path(source_root)
    source_file_path = normalize_mega_path(source_file)
    active_output_root = (
        normalize_mega_path(output_root)
        if output_root is not None
        else default_output_root(source_root_path)
    )

    relative_path = source_file_path.relative_to(source_root_path)
    output_relative_path = build_output_relative_path(relative_path, source_file_kind=source_file_kind)
    return active_output_root / output_relative_path


def build_output_relative_path(
    relative_path: PurePosixPath,
    *,
    source_file_kind: str,
) -> PurePosixPath:
    if source_file_kind == "media":
        return relative_path.with_suffix(".txt")
    if source_file_kind in {"image", "document"}:
        return append_name_before_suffix(relative_path, ".extracted.txt")
    raise ValueError(f"Unsupported source file kind for output path: {source_file_kind}")


def build_companion_output_relative_path(
    relative_path: PurePosixPath,
    *,
    source_file_kind: str,
    run_id: str,
) -> PurePosixPath:
    inserted_suffix = ".timestamps" if source_file_kind == "media" else ".structured"
    return append_name_before_suffix(relative_path, f"{inserted_suffix}.{run_id}.txt")


def build_companion_output_path(
    source_file: str | PurePosixPath,
    source_root: str | PurePosixPath,
    *,
    source_file_kind: str,
    run_id: str,
    output_root: str | PurePosixPath | None = None,
) -> PurePosixPath:
    source_root_path = normalize_mega_path(source_root)
    source_file_path = normalize_mega_path(source_file)
    active_output_root = (
        normalize_mega_path(output_root)
        if output_root is not None
        else default_output_root(source_root_path)
    )
    relative_path = source_file_path.relative_to(source_root_path)
    return active_output_root / build_companion_output_relative_path(
        relative_path,
        source_file_kind=source_file_kind,
        run_id=run_id,
    )


def filter_supported_media(paths: Iterable[str | PurePosixPath]) -> list[PurePosixPath]:
    supported = []
    for path in paths:
        normalized = normalize_mega_path(path)
        if normalized.suffix.lower() in SUPPORTED_MEDIA_SUFFIXES:
            supported.append(normalized)
    return supported


def filter_supported_local_media(paths: Iterable[Path]) -> list[Path]:
    supported = []
    for path in paths:
        candidate = Path(path)
        if candidate.is_file() and candidate.suffix.lower() in SUPPORTED_MEDIA_SUFFIXES:
            supported.append(candidate)
    return supported


def filter_processable_source_paths(paths: Iterable[str | PurePosixPath]) -> list[PurePosixPath]:
    normalized_paths = [normalize_mega_path(path) for path in paths]
    path_set = set(normalized_paths)
    processable_paths: list[PurePosixPath] = []
    for path in normalized_paths:
        if is_generated_pipeline_artifact(path):
            continue
        source_kind = classify_source_file_kind(path)
        if source_kind == "other":
            continue
        if media_transcript_source_for_path(path, path_set) is not None:
            continue
        processable_paths.append(path)
    return processable_paths


def filter_processable_local_sources(paths: Iterable[Path]) -> list[Path]:
    candidates = [Path(path) for path in paths if Path(path).is_file()]
    path_set = {Path(path) for path in candidates}
    processable: list[Path] = []
    for path in candidates:
        if is_generated_pipeline_artifact(path):
            continue
        source_kind = classify_source_file_kind(path)
        if source_kind == "other":
            continue
        if path.suffix.lower() == ".txt":
            for suffix in SUPPORTED_MEDIA_SUFFIXES:
                if path.with_suffix(suffix) in path_set:
                    break
            else:
                processable.append(path)
            continue
        processable.append(path)
    return processable


def is_mega_browser_folder_url(value: str) -> bool:
    return value.startswith("https://mega.nz/fm/") or value.startswith("https://mega.co.nz/fm/")


def is_mega_public_link(value: str) -> bool:
    return value.startswith("https://mega.nz/") or value.startswith("https://mega.co.nz/")


def resolve_source_spec(
    mega_source_path: str | None,
    mega_source_link: str | None,
    mega_output_path: str | None,
) -> SourceSpec:
    if bool(mega_source_path) == bool(mega_source_link):
        raise ValueError("Provide exactly one of --mega-source-path or --mega-source-link.")

    if mega_source_path is not None:
        if mega_output_path not in (None, ""):
            raise ValueError(
                "--mega-output-path is not supported when using --mega-source-path. "
                "Transcripts are written into the source folder."
            )
        source_root = normalize_mega_path(mega_source_path)
        return SourceSpec(
            source_kind="path",
            source_value=str(source_root),
            source_root=source_root,
            output_root=source_root,
            source_label=source_root.name,
        )

    assert mega_source_link is not None
    if is_mega_browser_folder_url(mega_source_link):
        raise ValueError(
            "Raw mega.nz/fm/ browser routes are not supported. Save a real MEGA path or public export link instead."
        )
    if not is_mega_public_link(mega_source_link):
        raise ValueError("MEGA source links must be public MEGA links.")
    if mega_output_path is None:
        raise ValueError("--mega-output-path is required when using --mega-source-link.")

    output_root = normalize_mega_path(mega_output_path)
    return SourceSpec(
        source_kind="link",
        source_value=mega_source_link,
        source_root=None,
        output_root=output_root,
        source_label=output_root.name,
    )


def build_s4_staging_key(
    source_label: str,
    run_id: str,
    relative_path: PurePosixPath,
    staging_prefix: str = "staging",
) -> str:
    return str(PurePosixPath(staging_prefix) / source_label / run_id / relative_path)


def build_s4_transcript_key(
    source_label: str,
    run_id: str,
    relative_path: PurePosixPath,
    transcript_prefix: str = "transcripts",
) -> str:
    return str(PurePosixPath(transcript_prefix) / source_label / run_id / relative_path)


def build_s4_companion_key(
    source_label: str,
    run_id: str,
    relative_path: PurePosixPath,
    transcript_prefix: str = "transcripts",
) -> str:
    return build_s4_transcript_key(
        source_label,
        run_id,
        relative_path,
        transcript_prefix=transcript_prefix,
    )


def build_s4_manifest_key(
    source_label: str,
    run_id: str,
    manifest_prefix: str = "manifests",
) -> str:
    return str(PurePosixPath(manifest_prefix) / source_label / f"{run_id}.json")


def build_local_artifact_root(work_dir: Path, source_label: str, run_id: str) -> Path:
    return Path(work_dir) / "exports" / source_label / run_id


def build_batch_runtime_item(
    batch_item: BatchItem,
    *,
    download_root: Path,
    artifact_root: Path,
    source_label: str,
    run_id: str,
    manifest_index: int,
    s4_staging_prefix: str,
    s4_transcript_prefix: str,
    s4_enabled: bool,
) -> BatchRuntimeItem:
    local_source_path = batch_item.local_source_path or (download_root / batch_item.relative_path)
    local_output_path = artifact_root / batch_item.output_relative_path
    local_companion_output_path = artifact_root / batch_item.companion_output_relative_path
    staging_key = None
    output_key = None
    companion_output_key = None
    if s4_enabled:
        staging_key = build_s4_staging_key(
            source_label,
            run_id,
            batch_item.relative_path,
            staging_prefix=s4_staging_prefix,
        )
        output_key = build_s4_transcript_key(
            source_label,
            run_id,
            batch_item.output_relative_path,
            transcript_prefix=s4_transcript_prefix,
        )
        companion_output_key = build_s4_companion_key(
            source_label,
            run_id,
            batch_item.companion_output_relative_path,
            transcript_prefix=s4_transcript_prefix,
        )
    return BatchRuntimeItem(
        batch_item=batch_item,
        local_source_path=local_source_path,
        local_output_path=local_output_path,
        local_companion_output_path=local_companion_output_path,
        staging_key=staging_key,
        output_key=output_key,
        companion_output_key=companion_output_key,
        manifest_index=manifest_index,
    )


def build_manifest_item(
    runtime_item: BatchRuntimeItem,
    *,
    status: str,
    resume_status: str | None = None,
    resume_reason: str | None = None,
    error: str | None = None,
    warnings: list[str] | None = None,
    processing_metadata: dict[str, str] | None = None,
) -> dict:
    manifest_item = {
        "source": runtime_item.batch_item.source_identifier,
        "source_file_kind": runtime_item.batch_item.source_file_kind,
        "output_path": str(runtime_item.batch_item.output_path),
        "companion_output_path": str(runtime_item.batch_item.companion_output_path),
        "status": status,
    }
    if resume_status is not None:
        manifest_item["resume_status"] = resume_status
    if resume_reason is not None:
        manifest_item["resume_reason"] = resume_reason
    if status != "skipped" and runtime_item.staging_key and runtime_item.output_key:
        manifest_item["s4_staging_key"] = runtime_item.staging_key
        manifest_item["s4_transcript_key"] = runtime_item.output_key
        if runtime_item.companion_output_key:
            manifest_item["s4_companion_key"] = runtime_item.companion_output_key
    if processing_metadata:
        manifest_item.update(processing_metadata)
    if error is not None:
        manifest_item["error"] = error
    if warnings:
        manifest_item["warnings"] = list(warnings)
    return manifest_item


def build_s4_warning(action: str, key: str, exc: Exception) -> str:
    return f"{action} failed for {key}: {exc}"


PREFETCH_POLL_INITIAL_SECONDS = 2
PREFETCH_POLL_MAX_SECONDS = 10
PREFETCH_TIMEOUT_SECONDS = 30 * 60  # 30 minutes


def wait_for_prefetch_ready(local_source_path: Path, timeout: float = PREFETCH_TIMEOUT_SECONDS) -> bool:
    """Poll for a `.ready` marker written by the prefetch process.

    Uses exponential backoff from 2s to 10s.  Returns True when the marker
    appears, False if the timeout expires.
    """
    marker = local_source_path.with_name(local_source_path.name + ".ready")
    deadline = time.monotonic() + timeout
    interval = PREFETCH_POLL_INITIAL_SECONDS
    while time.monotonic() < deadline:
        if marker.exists() and local_source_path.exists():
            return True
        time.sleep(min(interval, max(0, deadline - time.monotonic())))
        interval = min(interval * 1.5, PREFETCH_POLL_MAX_SECONDS)
    return marker.exists() and local_source_path.exists()


PREFETCH_MANIFEST_FILENAME = "_manifest.json"
PREFETCH_MANIFEST_TIMEOUT_SECONDS = 10 * 60  # 10 minutes


def read_prefetch_manifest(
    prefetch_root: Path,
    source_label: str,
    timeout: float = PREFETCH_MANIFEST_TIMEOUT_SECONDS,
) -> list[PurePosixPath] | None:
    """Read the file listing manifest written by the prefetch process.

    Polls for the manifest file with exponential backoff.  Returns
    a list of MEGA file paths, or None if the manifest doesn't appear
    within the timeout.
    """
    manifest_path = prefetch_root / source_label / PREFETCH_MANIFEST_FILENAME
    deadline = time.monotonic() + timeout
    interval = PREFETCH_POLL_INITIAL_SECONDS
    while time.monotonic() < deadline:
        if manifest_path.exists():
            try:
                data = json.loads(manifest_path.read_text())
                files = [PurePosixPath(f) for f in data.get("files", [])]
                print(
                    f"[prefetch] Read manifest for {source_label}: {len(files)} files",
                    file=sys.stderr,
                    flush=True,
                )
                return files
            except (json.JSONDecodeError, KeyError, OSError) as exc:
                print(
                    f"[prefetch] Failed to read manifest {manifest_path}: {exc}",
                    file=sys.stderr,
                    flush=True,
                )
                return None
        print(
            f"[prefetch] Waiting for manifest: {manifest_path}",
            file=sys.stderr,
            flush=True,
        )
        time.sleep(min(interval, max(0, deadline - time.monotonic())))
        interval = min(interval * 1.5, PREFETCH_POLL_MAX_SECONDS)
    print(
        f"[prefetch] Manifest not found after {timeout}s: {manifest_path}",
        file=sys.stderr,
        flush=True,
    )
    return None


def prepare_batch_item_io(
    mega_client,
    s4_client,
    runtime_item: BatchRuntimeItem,
    *,
    s4_bucket: str | None,
    prefetch_enabled: bool = False,
) -> list[str]:
    warnings: list[str] = []
    if runtime_item.batch_item.mega_source_path is not None:
        if prefetch_enabled:
            if not wait_for_prefetch_ready(runtime_item.local_source_path):
                warnings.append(
                    f"Prefetch timeout for {runtime_item.batch_item.mega_source_path}, falling back to direct download"
                )
                mega_client.download_file(runtime_item.batch_item.mega_source_path, runtime_item.local_source_path)
        else:
            mega_client.download_file(runtime_item.batch_item.mega_source_path, runtime_item.local_source_path)
    if s4_client is not None and s4_bucket and runtime_item.staging_key:
        try:
            s4_client.upload_file(runtime_item.local_source_path, s4_bucket, runtime_item.staging_key)
        except Exception as exc:
            warnings.append(build_s4_warning("s4 staging upload", runtime_item.staging_key, exc))
    return warnings


def finalize_batch_item_io(
    mega_client,
    s4_client,
    runtime_item: BatchRuntimeItem,
    *,
    s4_bucket: str | None,
) -> list[str]:
    warnings: list[str] = []
    mega_client.upload_file(runtime_item.local_output_path, runtime_item.batch_item.output_path)
    mega_client.upload_file(
        runtime_item.local_companion_output_path,
        runtime_item.batch_item.companion_output_path,
    )
    if s4_client is not None and s4_bucket and runtime_item.output_key and runtime_item.staging_key:
        try:
            s4_client.upload_file(runtime_item.local_output_path, s4_bucket, runtime_item.output_key)
        except Exception as exc:
            warnings.append(build_s4_warning("s4 transcript upload", runtime_item.output_key, exc))
        if runtime_item.companion_output_key:
            try:
                s4_client.upload_file(
                    runtime_item.local_companion_output_path,
                    s4_bucket,
                    runtime_item.companion_output_key,
                )
            except Exception as exc:
                warnings.append(
                    build_s4_warning("s4 companion upload", runtime_item.companion_output_key, exc)
                )
        try:
            s4_client.delete_object(s4_bucket, runtime_item.staging_key)
        except Exception as exc:
            warnings.append(build_s4_warning("s4 staging cleanup", runtime_item.staging_key, exc))
    return warnings


def discover_downloaded_source_root(local_dir: Path) -> Path:
    entries = [entry for entry in local_dir.iterdir() if not entry.name.startswith(".")]
    if len(entries) == 1 and entries[0].is_dir():
        return entries[0]
    return local_dir


def extract_error_text(result: subprocess.CompletedProcess[str], default_message: str) -> str:
    return result.stderr.strip() or result.stdout.strip() or default_message


def mega_download_staging_root() -> Path:
    configured_root = os.environ.get("MEGA_DOWNLOAD_TMPDIR")
    if configured_root:
        staging_root = Path(configured_root).expanduser()
    else:
        # Snap-packaged MEGAcmd cannot write into hidden home paths like ~/.cache.
        staging_root = Path.home() / "transcription-pipeline-downloads" / "mega-downloads"
    # On multi-GPU hosts, isolate staging dirs per process to prevent mega-cmd
    # temp file conflicts when TemporaryDirectory cleanup races with the daemon.
    cuda_device = os.environ.get("CUDA_VISIBLE_DEVICES")
    if cuda_device is not None:
        staging_root = staging_root / f"gpu{cuda_device}"
    staging_root.mkdir(parents=True, exist_ok=True)
    return staging_root


def mega_upload_staging_root() -> Path:
    configured_root = os.environ.get("MEGA_UPLOAD_TMPDIR")
    if configured_root:
        staging_root = Path(configured_root).expanduser()
    else:
        # Snap-packaged MEGAcmd can only read uploads from user-home paths.
        staging_root = Path.home() / "transcription-pipeline-uploads" / "mega-uploads"
    staging_root.mkdir(parents=True, exist_ok=True)
    return staging_root


def normalize_mega_ls_timestamp(value: str) -> str:
    candidate = value.strip()
    if not candidate:
        raise ValueError("MEGA timestamp cannot be empty.")

    if candidate.endswith("Z"):
        try:
            normalized = datetime.strptime(candidate, ISO_UTC_FORMAT)
            return normalized.strftime(ISO_UTC_FORMAT)
        except ValueError:
            pass

    for fmt in ("%Y-%m-%dT%H:%M:%S", SHORT_UTC_FORMAT):
        try:
            normalized = datetime.strptime(candidate, fmt)
            return normalized.strftime(ISO_UTC_FORMAT)
        except ValueError:
            continue

    raise ValueError(f"Unsupported MEGA timestamp format: {value}")


def parse_mega_ls_long_listing(
    output: str,
    root_path: str | PurePosixPath,
    *,
    file_time_field: str = "modified_at_utc",
) -> list[MegaLsEntry]:
    if file_time_field not in {"created_at_utc", "modified_at_utc"}:
        raise ValueError("file_time_field must be 'created_at_utc' or 'modified_at_utc'.")

    active_root = normalize_mega_path(root_path)
    current_dir = active_root
    entries: list[MegaLsEntry] = []

    for raw_line in output.splitlines():
        cleaned_line = raw_line.replace("\x00", "").rstrip()
        stripped = cleaned_line.strip()
        if not stripped:
            continue
        if stripped.startswith("[") or stripped.startswith("Resuming session"):
            continue
        if stripped.startswith("TRANSFERRING") or stripped.startswith("Download finished:"):
            continue
        if stripped.startswith("FLAGS "):
            continue
        if stripped.startswith("/") and stripped.endswith(":"):
            current_dir = normalize_mega_path(stripped[:-1])
            continue
        if stripped.endswith(":"):
            continue

        match = MEGA_LONG_ROW_RE.match(stripped)
        if match is None:
            continue

        kind = "folder" if match.group("flags").startswith("d") else "file"
        handle = match.group("handle")
        timestamp_utc = normalize_mega_ls_timestamp(match.group("date"))
        created_at_utc = timestamp_utc if kind == "folder" or file_time_field == "created_at_utc" else None
        modified_at_utc = timestamp_utc if kind == "file" and file_time_field == "modified_at_utc" else None
        entries.append(
            MegaLsEntry(
                path=current_dir / match.group("name"),
                parent_path=current_dir,
                name=match.group("name"),
                kind=kind,
                handle=handle[len(MEGA_HANDLE_PREFIX):] if handle.startswith(MEGA_HANDLE_PREFIX) else handle,
                created_at_utc=created_at_utc,
                modified_at_utc=modified_at_utc,
            )
        )

    return entries


def build_done_source_path(source_root: str | PurePosixPath) -> PurePosixPath:
    source_root_path = normalize_mega_path(source_root)
    if source_root_path.name.endswith(MEGA_DONE_SUFFIX):
        return source_root_path
    return source_root_path.parent / f"{source_root_path.name}{MEGA_DONE_SUFFIX}"


def rewrite_mega_path_prefix(
    remote_path: str,
    old_root: PurePosixPath,
    new_root: PurePosixPath,
) -> str:
    target = normalize_mega_path(remote_path)
    if target == old_root:
        return str(new_root)
    if str(target).startswith(f"{old_root}/"):
        return str(new_root / target.relative_to(old_root))
    return remote_path


class MegaCli:
    def __init__(self):
        self._command_cache: dict[str, str] = {}
        max_downloads = int(os.environ.get("MEGA_MAX_CONCURRENT_DOWNLOADS", "4"))
        self._download_semaphore = threading.Semaphore(max_downloads)

    def download_staging_root(self) -> Path:
        return mega_download_staging_root()

    def upload_staging_root(self) -> Path:
        return mega_upload_staging_root()

    def resolve_command(self, command_name: str) -> str:
        if command_name not in self._command_cache:
            candidate = which(command_name) or which(f"mega-cmd.{command_name}")
            if candidate is None:
                raise RuntimeError(f"Required MEGA command is not available: {command_name}")
            self._command_cache[command_name] = candidate
        return self._command_cache[command_name]

    def run(self, *args: str) -> subprocess.CompletedProcess[str]:
        return subprocess.run(
            args,
            capture_output=True,
            text=True,
            check=False,
        )

    def run_transfer(self, *args: str) -> subprocess.CompletedProcess[str]:
        result = subprocess.run(
            args,
            capture_output=True,
            text=True,
            check=False,
        )
        if result.returncode != 0 and not result.stdout.strip() and not result.stderr.strip():
            error_text = self.describe_error_code(result.returncode)
            return subprocess.CompletedProcess(
                args=args,
                returncode=result.returncode,
                stdout=error_text,
                stderr="",
            )
        return result

    def describe_error_code(self, returncode: int) -> str:
        result = self.run(
            self.resolve_command("mega-errorcode"),
            str(returncode),
        )
        return extract_error_text(result, f"MEGA transfer failed with exit code {returncode}")

    def list_nodes(self, source_root: str | PurePosixPath) -> list[MegaLsEntry]:
        source_root_path = normalize_mega_path(source_root)
        result = self.run(
            self.resolve_command("mega-ls"),
            "-l",
            "-R",
            "--show-handles",
            f"--time-format={MEGA_TIME_FORMAT}",
            str(source_root_path),
        )
        if result.returncode != 0:
            raise RuntimeError(extract_error_text(result, "mega-ls failed"))
        return parse_mega_ls_long_listing(result.stdout, source_root_path)

    def list_files(self, source_root: str | PurePosixPath) -> list[PurePosixPath]:
        return [entry.path for entry in self.list_nodes(source_root) if entry.kind == "file"]

    def resolve_handle_path(self, handle: str) -> PurePosixPath:
        result = self.run(self.resolve_command("mega-ls"), f"{MEGA_HANDLE_PREFIX}{handle}")
        if result.returncode != 0:
            raise RuntimeError(extract_error_text(result, "mega-ls failed"))
        for raw_line in result.stdout.splitlines():
            stripped = raw_line.strip()
            if stripped.startswith("/") and stripped.endswith(":"):
                return normalize_mega_path(stripped[:-1])
        raise RuntimeError(f"Unable to resolve folder handle: {handle}")

    def path_exists(self, remote_path: str | PurePosixPath) -> bool:
        result = self.run(self.resolve_command("mega-ls"), str(normalize_mega_path(remote_path)))
        return result.returncode == 0

    def ensure_directory(self, remote_path: str | PurePosixPath) -> None:
        target_path = normalize_mega_path(remote_path)
        result = self.run(
            self.resolve_command("mega-mkdir"),
            "-p",
            str(target_path),
        )
        if result.returncode == 0:
            return
        if self.path_exists(target_path):
            return
        raise RuntimeError(extract_error_text(result, "mega-mkdir failed"))

    def download_file(self, remote_path: str | PurePosixPath, local_path: Path) -> None:
        target_path = normalize_mega_path(remote_path)
        local_path = Path(local_path)
        local_path.parent.mkdir(parents=True, exist_ok=True)
        staging_root = self.download_staging_root()
        last_error = None
        with self._download_semaphore:
            for attempt in range(1, MEGA_DOWNLOAD_RETRY_ATTEMPTS + 1):
                with TemporaryDirectory(dir=str(staging_root)) as temp_dir:
                    temp_dir_path = Path(temp_dir)
                    result = self.run_transfer(
                        self.resolve_command("mega-get"),
                        str(target_path),
                        str(temp_dir_path),
                    )
                    if result.returncode == 0:
                        expected_path = temp_dir_path / target_path.name
                        if expected_path.exists():
                            shutil.move(str(expected_path), str(local_path))
                            return

                        downloaded_files = [entry for entry in temp_dir_path.iterdir() if entry.is_file()]
                        if len(downloaded_files) == 1:
                            shutil.move(str(downloaded_files[0]), str(local_path))
                            return

                        last_error = f"Downloaded file was not found for {target_path} temp_dir={temp_dir_path}"
                    else:
                        raw_stdout = result.stdout.strip() or "<empty>"
                        raw_stderr = result.stderr.strip() or "<empty>"
                        last_error = (
                            f"{extract_error_text(result, 'mega-get failed')} "
                            f"target={target_path} temp_dir={temp_dir_path} "
                            f"raw_stdout={raw_stdout!r} raw_stderr={raw_stderr!r}"
                        )

                if attempt >= MEGA_DOWNLOAD_RETRY_ATTEMPTS:
                    break

                # Cold-started MEGAcmd sessions can transiently fail the first
                # transfer; issue a lightweight listing before retrying.
                self.run(
                    self.resolve_command("mega-ls"),
                    "/",
                )
                time.sleep(MEGA_DOWNLOAD_RETRY_BACKOFF_SECONDS)

        raise RuntimeError(last_error or f"mega-get failed for {target_path}")

    def download_link(self, link: str, local_dir: Path) -> Path:
        local_dir.mkdir(parents=True, exist_ok=True)
        staging_root = self.download_staging_root()
        with TemporaryDirectory(dir=str(staging_root)) as temp_dir:
            temp_dir_path = Path(temp_dir)
            result = self.run_transfer(
                self.resolve_command("mega-get"),
                link,
                str(temp_dir_path),
            )
            if result.returncode != 0:
                raise RuntimeError(extract_error_text(result, "mega-get failed for link download"))

            downloaded_root = discover_downloaded_source_root(temp_dir_path)
            destination_root = local_dir / downloaded_root.name
            if destination_root.exists():
                if destination_root.is_dir():
                    shutil.rmtree(destination_root)
                else:
                    destination_root.unlink()
            shutil.move(str(downloaded_root), str(destination_root))
            return destination_root

    def upload_file(self, local_path: Path, remote_path: str | PurePosixPath) -> None:
        target_path = normalize_mega_path(remote_path)
        self.ensure_directory(target_path.parent)
        local_path = Path(local_path)
        staging_root = self.upload_staging_root()
        with TemporaryDirectory(dir=str(staging_root)) as temp_dir:
            temp_dir_path = Path(temp_dir)
            staged_path = temp_dir_path / local_path.name
            shutil.copy2(local_path, staged_path)
            result = self.run(
                self.resolve_command("mega-put"),
                str(staged_path),
                str(target_path.parent),
            )
        if result.returncode == 0:
            return
        if self.path_exists(target_path):
            return
        raise RuntimeError(extract_error_text(result, "mega-put failed"))

    def rename_path(self, old_path: str | PurePosixPath, new_path: str | PurePosixPath) -> None:
        current_path = normalize_mega_path(old_path)
        target_path = normalize_mega_path(new_path)
        if current_path == target_path:
            return
        result = self.run(
            self.resolve_command("mega-mv"),
            str(current_path),
            str(target_path),
        )
        if result.returncode == 0:
            return
        if self.path_exists(target_path) and not self.path_exists(current_path):
            return
        raise RuntimeError(extract_error_text(result, "mega-mv failed"))


class S4Client:
    def __init__(self, endpoint_url: str, access_key: str, secret_key: str):
        import boto3

        session = boto3.session.Session(
            aws_access_key_id=access_key,
            aws_secret_access_key=secret_key,
        )
        self._client = session.client("s3", endpoint_url=endpoint_url)

    @classmethod
    def from_env(cls, endpoint_url: str) -> "S4Client":
        access_key = os.environ.get("MEGA_S4_ACCESS_KEY")
        secret_key = os.environ.get("MEGA_S4_SECRET_KEY")
        if not access_key or not secret_key:
            raise RuntimeError("MEGA_S4_ACCESS_KEY and MEGA_S4_SECRET_KEY must be set.")
        return cls(
            endpoint_url=endpoint_url,
            access_key=access_key,
            secret_key=secret_key,
        )

    def upload_file(self, local_path: Path, bucket: str, key: str) -> None:
        self._client.upload_file(str(local_path), bucket, key)

    def delete_object(self, bucket: str, key: str) -> None:
        self._client.delete_object(Bucket=bucket, Key=key)


def format_timestamp_seconds(seconds: float) -> str:
    total_milliseconds = max(int(round(seconds * 1000)), 0)
    hours, remainder = divmod(total_milliseconds, 3_600_000)
    minutes, remainder = divmod(remainder, 60_000)
    secs, milliseconds = divmod(remainder, 1000)
    return f"{hours:02d}:{minutes:02d}:{secs:02d}.{milliseconds:03d}"


def write_text_artifact(path: Path, header_lines: list[str], body_lines: list[str]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8") as handle:
        if header_lines:
            handle.write("\n".join(header_lines).rstrip() + "\n\n")
        if body_lines:
            handle.write("\n".join(line.rstrip() for line in body_lines).rstrip() + "\n")


def build_media_output_lines(
    segments: list[dict[str, float | str]],
) -> tuple[list[str], list[str]]:
    canonical_lines = []
    companion_lines = []
    for segment in segments:
        text = str(segment.get("text", "")).strip()
        if not text:
            continue
        canonical_lines.append(text)
        companion_lines.append(
            f"[{format_timestamp_seconds(float(segment.get('start', 0.0)))}"
            f" --> {format_timestamp_seconds(float(segment.get('end', 0.0)))}] {text}"
        )
    return canonical_lines, companion_lines


def collapse_whitespace_text(value: str) -> str:
    return "\n".join(line.rstrip() for line in value.splitlines()).strip()


def strip_html_tags(value: str) -> str:
    without_blocks = re.sub(
        r"<(script|style)\b[^>]*>.*?</\1>",
        " ",
        value,
        flags=re.IGNORECASE | re.DOTALL,
    )
    without_tags = re.sub(r"<[^>]+>", " ", without_blocks)
    return html.unescape(re.sub(r"[ \t]+", " ", without_tags))


def structured_lines_from_plain_text(text: str, *, prefix: str = "line") -> list[str]:
    return [
        f"{prefix} {index:03d}: {line.strip()}"
        for index, line in enumerate(text.splitlines(), start=1)
        if line.strip()
    ]


def is_useful_text(text: str) -> bool:
    stripped = text.strip()
    return bool(stripped) and len(stripped) >= 8


class FasterWhisperTranscriber:
    def __init__(self, device: str = "cuda", compute_type: str = "auto"):
        self.device = device
        self.compute_type = compute_type
        self._models: dict[str, object] = {}
        self._fallback_transcriber = None

    def clone(self) -> "FasterWhisperTranscriber":
        cloned = FasterWhisperTranscriber(
            device=self.device,
            compute_type=self.compute_type,
        )
        if self._fallback_transcriber is not None:
            cloned._fallback_transcriber = self._fallback_transcriber.clone()
        return cloned

    def _load_model(self, model_name: str):
        if model_name not in self._models:
            from faster_whisper import WhisperModel

            self._models[model_name] = WhisperModel(
                model_name,
                device=self.device,
                compute_type=self.compute_type,
            )
        return self._models[model_name]

    def _load_fallback_transcriber(self):
        if self._fallback_transcriber is None:
            self._fallback_transcriber = TorchWhisperTranscriber(device=self.device)
        return self._fallback_transcriber

    def transcribe_segments(self, source_path: Path, model: str) -> tuple[list[dict[str, float | str]], dict[str, str], str]:
        try:
            whisper_model = self._load_model(model)
            segments, info = whisper_model.transcribe(str(source_path))
        except Exception as exc:
            if self.device == "cuda" and "not compiled with CUDA support" in str(exc):
                print(
                    "WARNING: faster-whisper CUDA backend unavailable; falling back to torch whisper.",
                    file=sys.stderr,
                )
                fallback_transcriber = self._load_fallback_transcriber()
                return fallback_transcriber.transcribe_segments(source_path, model)
            raise

        normalized_segments = [
            {
                "start": float(segment.start),
                "end": float(segment.end),
                "text": segment.text.strip(),
            }
            for segment in segments
            if segment.text.strip()
        ]
        metadata = {
            "language": info.language,
            "language_probability": f"{info.language_probability:.4f}",
        }
        return normalized_segments, metadata, "faster_whisper"

    def transcribe_file(self, source_path: Path, transcript_path: Path, model: str) -> None:
        segments, metadata, _ = self.transcribe_segments(source_path, model)
        canonical_lines, _ = build_media_output_lines(segments)
        header_lines = [f"# language={metadata['language']}"]
        if "language_probability" in metadata:
            header_lines.append(f"# probability={metadata['language_probability']}")
        write_text_artifact(
            transcript_path,
            header_lines,
            canonical_lines,
        )

    def transcribe_with_companion(
        self,
        source_path: Path,
        transcript_path: Path,
        companion_output_path: Path,
        *,
        model: str,
    ) -> ProcessResult:
        segments, metadata, processor_name = self.transcribe_segments(source_path, model)
        canonical_lines, companion_lines = build_media_output_lines(segments)
        header_lines = [f"# language={metadata['language']}"]
        if "language_probability" in metadata:
            header_lines.append(f"# probability={metadata['language_probability']}")
        write_text_artifact(transcript_path, header_lines, canonical_lines)
        write_text_artifact(
            companion_output_path,
            header_lines + [f"# content_mode={TIMESTAMP_CONTENT_MODE}"],
            companion_lines,
        )
        return ProcessResult(processor=processor_name, content_mode=TIMESTAMP_CONTENT_MODE)


class TorchWhisperTranscriber:
    def __init__(self, device: str = "cuda"):
        self.device = device
        self._models: dict[str, object] = {}

    def clone(self) -> "TorchWhisperTranscriber":
        return TorchWhisperTranscriber(device=self.device)

    def _load_model(self, model_name: str):
        if model_name not in self._models:
            import whisper

            self._models[model_name] = whisper.load_model(model_name, device=self.device)
        return self._models[model_name]

    def transcribe_segments(self, source_path: Path, model: str) -> tuple[list[dict[str, float | str]], dict[str, str], str]:
        import torch

        whisper_model = self._load_model(model)
        transcribe_kwargs = {}
        if self.device == "cuda":
            transcribe_kwargs["fp16"] = torch.cuda.is_available()
        result = whisper_model.transcribe(str(source_path), **transcribe_kwargs)
        segments = [
            {
                "start": float(segment.get("start", 0.0)),
                "end": float(segment.get("end", 0.0)),
                "text": str(segment.get("text", "")).strip(),
            }
            for segment in result.get("segments", [])
            if str(segment.get("text", "")).strip()
        ]
        metadata = {
            "language": str(result.get("language") or "unknown"),
        }
        return segments, metadata, "torch_whisper"

    def transcribe_file(self, source_path: Path, transcript_path: Path, model: str) -> None:
        segments, metadata, _ = self.transcribe_segments(source_path, model)
        canonical_lines, _ = build_media_output_lines(segments)
        write_text_artifact(
            transcript_path,
            [f"# language={metadata['language']}"],
            canonical_lines,
        )

    def transcribe_with_companion(
        self,
        source_path: Path,
        transcript_path: Path,
        companion_output_path: Path,
        *,
        model: str,
    ) -> ProcessResult:
        segments, metadata, processor_name = self.transcribe_segments(source_path, model)
        canonical_lines, companion_lines = build_media_output_lines(segments)
        header_lines = [f"# language={metadata['language']}"]
        write_text_artifact(transcript_path, header_lines, canonical_lines)
        write_text_artifact(
            companion_output_path,
            header_lines + [f"# content_mode={TIMESTAMP_CONTENT_MODE}"],
            companion_lines,
        )
        return ProcessResult(processor=processor_name, content_mode=TIMESTAMP_CONTENT_MODE)


class DocumentTextExtractor:
    def clone(self) -> "DocumentTextExtractor":
        return DocumentTextExtractor()

    def extract_with_companion(
        self,
        source_path: Path,
        transcript_path: Path,
        companion_output_path: Path,
        *,
        source_file_kind: str,
    ) -> ProcessResult:
        suffix = source_path.suffix.lower()
        native_error = None
        if suffix in SUPPORTED_NATIVE_DOCUMENT_SUFFIXES:
            try:
                text, structured_text, processor_name = self._extract_native(source_path)
            except Exception as exc:  # pragma: no cover - optional dependency path.
                native_error = exc
            else:
                if is_useful_text(text) or suffix in {".txt", ".md", ".json", ".html", ".htm", ".csv", ".tsv", ".rtf", ".srt"}:
                    self._write_outputs(
                        transcript_path,
                        companion_output_path,
                        text=text,
                        structured_text=structured_text,
                        processor_name=processor_name,
                        source_file_kind=source_file_kind,
                    )
                    return ProcessResult(processor=processor_name, content_mode=STRUCTURED_CONTENT_MODE)
        if source_file_kind in {"image", "document"} and (
            suffix in SUPPORTED_IMAGE_SUFFIXES or suffix in SUPPORTED_OCR_DOCUMENT_SUFFIXES or native_error is not None
        ):
            try:
                text, structured_text = self._extract_with_marker(source_path, force_ocr=source_file_kind == "image" or suffix == ".pdf")
            except Exception as exc:  # pragma: no cover - runtime dependency path.
                if native_error is not None:
                    raise RuntimeError(f"{native_error}; marker fallback failed: {exc}") from exc
                raise
            self._write_outputs(
                transcript_path,
                companion_output_path,
                text=text,
                structured_text=structured_text,
                processor_name=MARKER_PROCESSOR,
                source_file_kind=source_file_kind,
            )
            return ProcessResult(processor=MARKER_PROCESSOR, content_mode=STRUCTURED_CONTENT_MODE)
        if native_error is not None:
            raise RuntimeError(str(native_error))
        raise RuntimeError(f"Unsupported extractable file: {source_path}")

    def _write_outputs(
        self,
        transcript_path: Path,
        companion_output_path: Path,
        *,
        text: str,
        structured_text: str,
        processor_name: str,
        source_file_kind: str,
    ) -> None:
        write_text_artifact(
            transcript_path,
            [
                f"# processor={processor_name}",
                f"# source_kind={source_file_kind}",
            ],
            text.splitlines(),
        )
        write_text_artifact(
            companion_output_path,
            [
                f"# processor={processor_name}",
                f"# source_kind={source_file_kind}",
                f"# content_mode={STRUCTURED_CONTENT_MODE}",
            ],
            structured_text.splitlines(),
        )

    def _extract_native(self, source_path: Path) -> tuple[str, str, str]:
        suffix = source_path.suffix.lower()
        if suffix in {".txt", ".md", ".srt"}:
            text = source_path.read_text(encoding="utf-8")
            return collapse_whitespace_text(text), "\n".join(structured_lines_from_plain_text(text)), f"{NATIVE_PROCESSOR_PREFIX}_text"
        if suffix == ".json":
            loaded = json.loads(source_path.read_text(encoding="utf-8"))
            text = json.dumps(loaded, indent=2, ensure_ascii=False)
            return text, "\n".join(structured_lines_from_plain_text(text, prefix="json_line")), f"{NATIVE_PROCESSOR_PREFIX}_json"
        if suffix in {".html", ".htm"}:
            raw_html = source_path.read_text(encoding="utf-8")
            text = collapse_whitespace_text(strip_html_tags(raw_html))
            return text, "\n".join(structured_lines_from_plain_text(text)), f"{NATIVE_PROCESSOR_PREFIX}_html"
        if suffix in {".csv", ".tsv"}:
            delimiter = "," if suffix == ".csv" else "\t"
            with source_path.open("r", encoding="utf-8", newline="") as handle:
                rows = [row for row in csv.reader(handle, delimiter=delimiter)]
            row_lines = ["\t".join(cell.strip() for cell in row if cell.strip()) for row in rows]
            row_lines = [line for line in row_lines if line.strip()]
            structured = [f"row {index:03d}: {line}" for index, line in enumerate(row_lines, start=1)]
            return "\n".join(row_lines), "\n".join(structured), f"{NATIVE_PROCESSOR_PREFIX}_table"
        if suffix in {".xlsx", ".xls"}:
            return self._extract_spreadsheet(source_path)
        if suffix == ".docx":
            return self._extract_docx(source_path)
        if suffix == ".pptx":
            return self._extract_pptx(source_path)
        if suffix == ".rtf":
            raw_rtf = source_path.read_text(encoding="utf-8", errors="ignore")
            text = collapse_whitespace_text(re.sub(r"\\[a-z]+\d* ?", " ", raw_rtf).replace("{", " ").replace("}", " "))
            return text, "\n".join(structured_lines_from_plain_text(text)), f"{NATIVE_PROCESSOR_PREFIX}_rtf"
        raise RuntimeError(f"No native extractor for {source_path.suffix}")

    def _extract_spreadsheet(self, source_path: Path) -> tuple[str, str, str]:
        if source_path.suffix.lower() == ".xlsx":
            try:
                from openpyxl import load_workbook
            except ImportError as exc:  # pragma: no cover - optional dependency path.
                raise RuntimeError("openpyxl is required for .xlsx extraction") from exc

            workbook = load_workbook(source_path, read_only=True, data_only=True)
            text_lines: list[str] = []
            structured_lines: list[str] = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_lines.append(f"# sheet: {sheet_name}")
                structured_lines.append(f"sheet: {sheet_name}")
                for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                    cells = [str(cell).strip() for cell in row if cell not in (None, "")]
                    if not cells:
                        continue
                    joined = "\t".join(cells)
                    text_lines.append(joined)
                    structured_lines.append(f"sheet {sheet_name} | row {row_index:03d}: {joined}")
            return "\n".join(text_lines), "\n".join(structured_lines), f"{NATIVE_PROCESSOR_PREFIX}_spreadsheet"
        try:
            import pandas as pd
        except ImportError as exc:  # pragma: no cover - optional dependency path.
            raise RuntimeError("pandas is required for .xls extraction") from exc

        sheets = pd.read_excel(source_path, sheet_name=None)
        text_lines: list[str] = []
        structured_lines: list[str] = []
        for sheet_name, dataframe in sheets.items():
            text_lines.append(f"# sheet: {sheet_name}")
            structured_lines.append(f"sheet: {sheet_name}")
            for row_index, row in enumerate(dataframe.fillna("").itertuples(index=False), start=1):
                cells = [str(cell).strip() for cell in row if str(cell).strip()]
                if not cells:
                    continue
                joined = "\t".join(cells)
                text_lines.append(joined)
                structured_lines.append(f"sheet {sheet_name} | row {row_index:03d}: {joined}")
        return "\n".join(text_lines), "\n".join(structured_lines), f"{NATIVE_PROCESSOR_PREFIX}_spreadsheet"

    def _extract_docx(self, source_path: Path) -> tuple[str, str, str]:
        try:
            from docx import Document
        except ImportError as exc:  # pragma: no cover - optional dependency path.
            raise RuntimeError("python-docx is required for .docx extraction") from exc

        document = Document(source_path)
        paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
        text = "\n\n".join(paragraphs)
        structured = "\n".join(
            f"paragraph {index:03d}: {paragraph}"
            for index, paragraph in enumerate(paragraphs, start=1)
        )
        return text, structured, f"{NATIVE_PROCESSOR_PREFIX}_docx"

    def _extract_pptx(self, source_path: Path) -> tuple[str, str, str]:
        try:
            from pptx import Presentation
        except ImportError as exc:  # pragma: no cover - optional dependency path.
            raise RuntimeError("python-pptx is required for .pptx extraction") from exc

        presentation = Presentation(source_path)
        text_lines: list[str] = []
        structured_lines: list[str] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    candidate = str(shape.text).strip()
                    if candidate:
                        slide_text.append(candidate)
            if not slide_text:
                continue
            text_lines.append(f"# slide {slide_index}")
            text_lines.extend(slide_text)
            structured_lines.append(f"slide {slide_index:03d}: {' | '.join(slide_text)}")
        return "\n".join(text_lines), "\n".join(structured_lines), f"{NATIVE_PROCESSOR_PREFIX}_pptx"

    def _extract_with_marker(self, source_path: Path, *, force_ocr: bool) -> tuple[str, str]:
        marker_command = which("marker_single")
        if marker_command is None:
            raise RuntimeError("marker_single is not installed or not on PATH.")
        with TemporaryDirectory() as temp_dir:
            output_dir = Path(temp_dir)
            command = [
                marker_command,
                str(source_path),
                "--output_dir",
                str(output_dir),
                "--output_format",
                "markdown",
                "--paginate_output",
            ]
            if force_ocr:
                command.append("--force_ocr")
            result = subprocess.run(command, capture_output=True, text=True, check=False)
            if result.returncode != 0:
                raise RuntimeError(extract_error_text(result, "marker_single failed"))
            markdown_files = sorted(output_dir.rglob("*.md"))
            if not markdown_files:
                raise RuntimeError(f"marker_single did not produce markdown for {source_path.name}")
            markdown_text = markdown_files[0].read_text(encoding="utf-8")
            canonical_text = re.sub(r"\n\s*\n\d+\n-+\n\s*\n", "\n\n", markdown_text).strip()
            structured_text = markdown_text.strip()
            return canonical_text, structured_text


class MixedFileProcessor:
    def __init__(self, media_transcriber, document_extractor: DocumentTextExtractor | None = None):
        self.media_transcriber = media_transcriber
        self.document_extractor = document_extractor or DocumentTextExtractor()

    def clone(self) -> "MixedFileProcessor":
        media_clone = self.media_transcriber.clone() if hasattr(self.media_transcriber, "clone") else self.media_transcriber
        return MixedFileProcessor(
            media_transcriber=media_clone,
            document_extractor=self.document_extractor.clone(),
        )

    def process_item(self, runtime_item: BatchRuntimeItem, *, model: str) -> ProcessResult:
        runtime_item.local_output_path.parent.mkdir(parents=True, exist_ok=True)
        runtime_item.local_companion_output_path.parent.mkdir(parents=True, exist_ok=True)
        source_kind = runtime_item.batch_item.source_file_kind
        if source_kind == "media":
            if hasattr(self.media_transcriber, "transcribe_with_companion"):
                return self.media_transcriber.transcribe_with_companion(
                    runtime_item.local_source_path,
                    runtime_item.local_output_path,
                    runtime_item.local_companion_output_path,
                    model=model,
                )
            self.media_transcriber.transcribe_file(
                runtime_item.local_source_path,
                runtime_item.local_output_path,
                model=model,
            )
            fallback_text = runtime_item.local_output_path.read_text(encoding="utf-8")
            write_text_artifact(
                runtime_item.local_companion_output_path,
                [f"# processor={MEDIA_PROCESSOR}", f"# content_mode={TIMESTAMP_CONTENT_MODE}"],
                structured_lines_from_plain_text(fallback_text, prefix="segment"),
            )
            return ProcessResult(processor=MEDIA_PROCESSOR, content_mode=TIMESTAMP_CONTENT_MODE)
        if source_kind in {"image", "document"}:
            return self.document_extractor.extract_with_companion(
                runtime_item.local_source_path,
                runtime_item.local_output_path,
                runtime_item.local_companion_output_path,
                source_file_kind=source_kind,
            )
        raise RuntimeError(f"Unsupported source file kind: {source_kind}")


def plan_batch(
    source_paths: Iterable[str | PurePosixPath],
    mega_source_path: str | PurePosixPath,
    *,
    run_id: str,
    mega_output_path: str | PurePosixPath | None = None,
) -> list[BatchItem]:
    source_root = normalize_mega_path(mega_source_path)
    output_root = (
        normalize_mega_path(mega_output_path)
        if mega_output_path is not None
        else default_output_root(source_root)
    )
    return [
        BatchItem(
            source_identifier=str(source_path),
            relative_path=relative_path,
            source_file_kind=source_file_kind,
            output_relative_path=build_output_relative_path(
                relative_path,
                source_file_kind=source_file_kind,
            ),
            companion_output_relative_path=build_companion_output_relative_path(
                relative_path,
                source_file_kind=source_file_kind,
                run_id=run_id,
            ),
            mega_source_path=normalize_mega_path(source_path),
            output_path=build_generated_output_path(
                source_path,
                source_root,
                source_file_kind=source_file_kind,
                output_root=output_root,
            ),
            companion_output_path=build_companion_output_path(
                source_path,
                source_root,
                source_file_kind=source_file_kind,
                run_id=run_id,
                output_root=output_root,
            ),
        )
        for source_path in filter_processable_source_paths(source_paths)
        for relative_path in [normalize_mega_path(source_path).relative_to(source_root)]
        for source_file_kind in [classify_source_file_kind(source_path)]
    ]


def write_manifest(manifest: dict, manifest_path: Path) -> None:
    manifest_path.parent.mkdir(parents=True, exist_ok=True)
    manifest_path.write_text(json.dumps(manifest, indent=2), encoding="utf-8")


def plan_link_batch(
    local_source_root: Path,
    output_root: PurePosixPath,
    source_link: str,
    *,
    run_id: str,
) -> list[BatchItem]:
    return [
        BatchItem(
            source_identifier=f"{source_link}::{local_path.relative_to(local_source_root).as_posix()}",
            relative_path=relative_path,
            source_file_kind=source_file_kind,
            output_relative_path=build_output_relative_path(
                relative_path,
                source_file_kind=source_file_kind,
            ),
            companion_output_relative_path=build_companion_output_relative_path(
                relative_path,
                source_file_kind=source_file_kind,
                run_id=run_id,
            ),
            local_source_path=local_path,
            output_path=output_root
            / build_output_relative_path(relative_path, source_file_kind=source_file_kind),
            companion_output_path=output_root
            / build_companion_output_relative_path(
                relative_path,
                source_file_kind=source_file_kind,
                run_id=run_id,
            ),
        )
        for local_path in filter_processable_local_sources(sorted(local_source_root.rglob("*")))
        for relative_path in [PurePosixPath(local_path.relative_to(local_source_root).as_posix())]
        for source_file_kind in [classify_source_file_kind(local_path)]
    ]


def build_processing_runner(processor, *, model: str, transcribe_workers: int):
    if transcribe_workers == 1:

        def run(runtime_item: BatchRuntimeItem) -> tuple[BatchRuntimeItem, ProcessResult]:
            return runtime_item, processor.process_item(runtime_item, model=model)

        return run

    clone_method = getattr(processor, "clone", None)
    if not callable(clone_method):
        raise RuntimeError(
            "Parallel transcription requires a transcriber with a clone() method."
        )

    worker_state = threading.local()
    worker_assignment_lock = threading.Lock()
    worker_transcribers = deque()

    for _ in range(transcribe_workers):
        worker_transcriber = clone_method()
        if worker_transcriber is None:
            raise RuntimeError("Transcriber clone() returned no worker instance.")
        worker_transcribers.append(worker_transcriber)

    def resolve_worker_transcriber():
        worker_transcriber = getattr(worker_state, "transcriber", None)
        if worker_transcriber is None:
            with worker_assignment_lock:
                if not worker_transcribers:
                    raise RuntimeError("No transcriber clones remain for worker assignment.")
                worker_transcriber = worker_transcribers.popleft()
            worker_state.transcriber = worker_transcriber
        return worker_transcriber

    def run(runtime_item: BatchRuntimeItem) -> tuple[BatchRuntimeItem, ProcessResult]:
        worker_transcriber = resolve_worker_transcriber()
        return runtime_item, worker_transcriber.process_item(runtime_item, model=model)

    return run


def run_batch(
    mega_client,
    s4_client,
    transcriber,
    mega_source_path: str | None,
    work_dir: Path,
    s4_bucket: str | None,
    mega_source_link: str | None = None,
    mega_output_path: str | None = None,
    s4_staging_prefix: str = "staging",
    s4_transcript_prefix: str = "transcripts",
    s4_manifest_prefix: str = "manifests",
    model: str = "large-v3",
    force: bool = False,
    io_workers: int = 4,
    transcribe_workers: int = 1,
    overlap_io_with_transcription: bool = False,
    timestamp_value: str | None = None,
    include_file_kinds: str | Iterable[str] | None = None,
    disable_source_rename: bool = False,
    compute_profile: str | None = None,
    selection_policy: str | None = None,
    preferred_regions: str | Iterable[str] | None = None,
    selected_region: str | None = None,
    instance_type: str | None = None,
    selected_image_id: str | None = None,
    image_family: str | None = None,
    image_version: str | None = None,
    architecture: str | None = None,
    gpu_count: int | None = None,
    resolved_runtime_root: str | None = None,
    attached_volume_mount: str | None = None,
    scheduler_batch_id: str | None = None,
    scheduler_policy: str | None = None,
    worker_name: str | None = None,
    runner_host: str | None = None,
    prefetch_root: str | None = None,
    file_slice: str | None = None,
) -> dict:
    source_spec = resolve_source_spec(mega_source_path, mega_source_link, mega_output_path)
    run_id = timestamp_value or utc_timestamp()
    processor = transcriber if hasattr(transcriber, "process_item") else MixedFileProcessor(transcriber)
    selected_file_kinds = normalize_selected_file_kinds(include_file_kinds)
    normalized_preferred_regions = normalize_preferred_regions(preferred_regions)
    s4_enabled = s4_client is not None and bool(s4_bucket)
    if io_workers < 1:
        raise ValueError("io_workers must be at least 1.")
    if transcribe_workers < 1:
        raise ValueError("transcribe_workers must be at least 1.")
    if transcribe_workers > MAX_TRANSCRIBE_WORKERS:
        raise ValueError(f"transcribe_workers={transcribe_workers} exceeds maximum of {MAX_TRANSCRIBE_WORKERS}.")
    if transcribe_workers > 1 and not hasattr(transcriber, "process_item") and not callable(getattr(transcriber, "clone", None)):
        raise RuntimeError("Parallel transcription requires a transcriber with a clone() method.")
    effective_overlap = overlap_io_with_transcription or transcribe_workers > 1
    prefetch_enabled = prefetch_root is not None
    prefetch_source_root = (
        Path(prefetch_root) / source_spec.source_label if prefetch_root else None
    )
    artifact_root = build_local_artifact_root(Path(work_dir), source_spec.source_label, run_id)
    download_root = Path(work_dir) / "downloads" / source_spec.source_label / run_id
    manifest_results: list[dict | None] = []
    active_runtime_items: list[BatchRuntimeItem] = []
    item_warnings: dict[int, list[str]] = {}
    item_processing_metadata: dict[int, dict[str, str]] = {}
    item_resume_metadata: dict[int, ResumeDecision] = {}
    processed = skipped = failed = 0
    rename_applied = False
    rename_error = None
    source_path_before = source_spec.source_root
    source_path_after = source_spec.source_root or source_spec.output_root
    manifest_warnings: list[str] = []
    resume_summary = {
        "completed_skipped": 0,
        "pending_queued": 0,
        "partial_requeued": 0,
        "failed_requeued": 0,
    }

    if source_spec.source_kind == "path":
        assert source_spec.source_root is not None
        # When prefetch is enabled, read the manifest written by the prefetch
        # process instead of calling mega-ls.  This avoids overloading the
        # single mega-cmd-server daemon on multi-GPU hosts.
        if prefetch_root:
            manifest_files = read_prefetch_manifest(
                Path(prefetch_root),
                source_spec.source_label,
            )
            if manifest_files is not None:
                source_paths = manifest_files
                print(f"[prefetch] Using manifest ({len(source_paths)} files) instead of mega-ls", file=sys.stderr, flush=True)
            else:
                print("[prefetch] Manifest unavailable, falling back to mega-ls", file=sys.stderr, flush=True)
                source_paths = mega_client.list_files(source_spec.source_root)
        else:
            source_paths = mega_client.list_files(source_spec.source_root)
        remote_resume_paths = build_remote_resume_paths(
            source_paths,
            root=source_spec.source_root,
        )
        batch = [
            item
            for item in plan_batch(
            source_paths,
            mega_source_path=source_spec.source_root,
            run_id=run_id,
            mega_output_path=source_spec.output_root,
            )
            if item.source_file_kind in selected_file_kinds
        ]
    else:
        try:
            remote_output_paths = mega_client.list_files(source_spec.output_root)
        except Exception:
            remote_output_paths = []
        remote_resume_paths = build_remote_resume_paths(
            remote_output_paths,
            root=source_spec.output_root,
        )
        shared_root = mega_client.download_link(
            source_spec.source_value,
            Path(work_dir) / "shared-links" / source_spec.source_label / run_id,
        )
        batch = [
            item
            for item in plan_link_batch(
                shared_root,
                source_spec.output_root,
                source_spec.source_value,
                run_id=run_id,
            )
            if item.source_file_kind in selected_file_kinds
        ]
        mega_client.ensure_directory(source_spec.output_root)

    # File-level slicing: distribute files across multiple GPU workers.
    if file_slice:
        slice_index, slice_total = (int(x) for x in file_slice.split("/"))
        pre_slice_count = len(batch)
        batch = [item for i, item in enumerate(batch) if i % slice_total == slice_index]
        print(f"[file-slice] Slice {slice_index}/{slice_total}: {len(batch)} of {pre_slice_count} files assigned to this worker", file=sys.stderr, flush=True)

    effective_download_root = prefetch_source_root if prefetch_enabled else download_root
    for item in batch:
        runtime_item = build_batch_runtime_item(
            item,
            download_root=effective_download_root,
            artifact_root=artifact_root,
            source_label=source_spec.source_label,
            run_id=run_id,
            manifest_index=len(manifest_results),
            s4_staging_prefix=s4_staging_prefix,
            s4_transcript_prefix=s4_transcript_prefix,
            s4_enabled=s4_enabled,
        )
        resume_decision = classify_resume_decision(
            item,
            remote_resume_paths,
            force=force,
        )
        item_resume_metadata[runtime_item.manifest_index] = resume_decision
        update_resume_summary(resume_summary, resume_decision)
        if not resume_decision.queue_for_processing:
            skipped += 1
            manifest_results.append(
                build_manifest_item(
                    runtime_item,
                    status="skipped",
                    resume_status=resume_decision.status,
                    resume_reason=resume_decision.reason,
                    processing_metadata={
                        "processor": "existing_output",
                        "content_mode": (
                            TIMESTAMP_CONTENT_MODE if item.source_file_kind == "media" else STRUCTURED_CONTENT_MODE
                        ),
                    },
                )
            )
            continue
        manifest_results.append(None)
        active_runtime_items.append(runtime_item)

    # When prefetch is enabled, reorder so files with existing .ready markers
    # are processed first.  This prevents the IO pool from filling up with
    # threads blocked on wait_for_prefetch_ready while ready files sit idle.
    if active_runtime_items and prefetch_enabled:
        def _ready_sort_key(item: BatchRuntimeItem) -> int:
            marker = item.local_source_path.with_name(item.local_source_path.name + ".ready")
            return 0 if marker.exists() else 1
        active_runtime_items.sort(key=_ready_sort_key)
        ready_count = sum(1 for item in active_runtime_items if item.local_source_path.with_name(item.local_source_path.name + ".ready").exists())
        print(f"[prefetch] Reordered batch: {ready_count} ready first, {len(active_runtime_items) - ready_count} waiting", file=sys.stderr, flush=True)

    if active_runtime_items:
        transcribe_runtime_item = build_processing_runner(
            processor,
            model=model,
            transcribe_workers=transcribe_workers,
        )
        if not effective_overlap:
            for runtime_item in active_runtime_items:
                try:
                    prepare_warnings = prepare_batch_item_io(
                        mega_client,
                        s4_client,
                        runtime_item,
                        s4_bucket=s4_bucket,
                        prefetch_enabled=prefetch_enabled,
                    )
                    if prepare_warnings:
                        item_warnings.setdefault(runtime_item.manifest_index, []).extend(prepare_warnings)
                    _, process_result = transcribe_runtime_item(runtime_item)
                    item_processing_metadata[runtime_item.manifest_index] = {
                        "processor": process_result.processor,
                        "content_mode": process_result.content_mode,
                    }
                    publish_warnings = finalize_batch_item_io(
                        mega_client,
                        s4_client,
                        runtime_item,
                        s4_bucket=s4_bucket,
                    )
                    if publish_warnings:
                        item_warnings.setdefault(runtime_item.manifest_index, []).extend(publish_warnings)
                except Exception as exc:  # pragma: no cover - exercised in tests via fakes.
                    failed += 1
                    manifest_results[runtime_item.manifest_index] = build_manifest_item(
                        runtime_item,
                        status="failed",
                        resume_status=item_resume_metadata[runtime_item.manifest_index].status,
                        resume_reason=item_resume_metadata[runtime_item.manifest_index].reason,
                        error=str(exc),
                        warnings=item_warnings.get(runtime_item.manifest_index),
                        processing_metadata=item_processing_metadata.get(runtime_item.manifest_index),
                    )
                else:
                    processed += 1
                    manifest_results[runtime_item.manifest_index] = build_manifest_item(
                        runtime_item,
                        status="processed",
                        resume_status=item_resume_metadata[runtime_item.manifest_index].status,
                        resume_reason=item_resume_metadata[runtime_item.manifest_index].reason,
                        warnings=item_warnings.get(runtime_item.manifest_index),
                        processing_metadata=item_processing_metadata.get(runtime_item.manifest_index),
                    )
        else:
            with (
                ThreadPoolExecutor(max_workers=io_workers) as io_pool,
                ThreadPoolExecutor(max_workers=transcribe_workers) as transcribe_pool,
            ):
                pending_prepares: dict[object, BatchRuntimeItem] = {}
                pending_transcribes: dict[object, BatchRuntimeItem] = {}
                pending_publishes: dict[object, BatchRuntimeItem] = {}
                ready_runtime_items: deque[BatchRuntimeItem] = deque()
                runtime_iter = iter(active_runtime_items)

                def queue_prepare() -> bool:
                    try:
                        next_item = next(runtime_iter)
                    except StopIteration:
                        return False
                    prepare_future = io_pool.submit(
                        prepare_batch_item_io,
                        mega_client,
                        s4_client,
                        next_item,
                        s4_bucket=s4_bucket,
                        prefetch_enabled=prefetch_enabled,
                    )
                    pending_prepares[prepare_future] = next_item
                    return True

                def fill_prepare_queue() -> None:
                    while len(pending_prepares) < io_workers and queue_prepare():
                        pass

                fill_prepare_queue()

                while pending_prepares or pending_transcribes or pending_publishes or ready_runtime_items:
                    while ready_runtime_items and len(pending_transcribes) < transcribe_workers:
                        runtime_item = ready_runtime_items.popleft()
                        transcribe_future = transcribe_pool.submit(
                            transcribe_runtime_item,
                            runtime_item,
                        )
                        pending_transcribes[transcribe_future] = runtime_item

                    active_futures = [
                        *pending_prepares.keys(),
                        *pending_transcribes.keys(),
                        *pending_publishes.keys(),
                    ]
                    if not active_futures:
                        continue

                    done_futures, _ = wait(active_futures, return_when=FIRST_COMPLETED)
                    for completed_future in done_futures:
                        runtime_item = pending_prepares.pop(completed_future, None)
                        if runtime_item is not None:
                            try:
                                prepare_warnings = completed_future.result()
                            except Exception as exc:  # pragma: no cover - exercised in tests via fakes.
                                failed += 1
                                manifest_results[runtime_item.manifest_index] = build_manifest_item(
                                    runtime_item,
                                    status="failed",
                                    resume_status=item_resume_metadata[runtime_item.manifest_index].status,
                                    resume_reason=item_resume_metadata[runtime_item.manifest_index].reason,
                                    error=str(exc),
                                    warnings=item_warnings.get(runtime_item.manifest_index),
                                )
                            else:
                                if prepare_warnings:
                                    item_warnings.setdefault(runtime_item.manifest_index, []).extend(prepare_warnings)
                                ready_runtime_items.append(runtime_item)
                            continue

                        runtime_item = pending_transcribes.pop(completed_future, None)
                        if runtime_item is not None:
                            try:
                                _, process_result = completed_future.result()
                            except Exception as exc:  # pragma: no cover - exercised in tests via fakes.
                                failed += 1
                                manifest_results[runtime_item.manifest_index] = build_manifest_item(
                                    runtime_item,
                                    status="failed",
                                    resume_status=item_resume_metadata[runtime_item.manifest_index].status,
                                    resume_reason=item_resume_metadata[runtime_item.manifest_index].reason,
                                    error=str(exc),
                                    warnings=item_warnings.get(runtime_item.manifest_index),
                                    processing_metadata=item_processing_metadata.get(runtime_item.manifest_index),
                                )
                            else:
                                item_processing_metadata[runtime_item.manifest_index] = {
                                    "processor": process_result.processor,
                                    "content_mode": process_result.content_mode,
                                }
                                publish_future = io_pool.submit(
                                    finalize_batch_item_io,
                                    mega_client,
                                    s4_client,
                                    runtime_item,
                                    s4_bucket=s4_bucket,
                                )
                                pending_publishes[publish_future] = runtime_item
                            continue

                        runtime_item = pending_publishes.pop(completed_future)
                        try:
                            publish_warnings = completed_future.result()
                        except Exception as exc:  # pragma: no cover - exercised in tests via fakes.
                            failed += 1
                            manifest_results[runtime_item.manifest_index] = build_manifest_item(
                                runtime_item,
                                status="failed",
                                resume_status=item_resume_metadata[runtime_item.manifest_index].status,
                                resume_reason=item_resume_metadata[runtime_item.manifest_index].reason,
                                error=str(exc),
                                warnings=item_warnings.get(runtime_item.manifest_index),
                                processing_metadata=item_processing_metadata.get(runtime_item.manifest_index),
                            )
                        else:
                            if publish_warnings:
                                item_warnings.setdefault(runtime_item.manifest_index, []).extend(publish_warnings)
                            processed += 1
                            manifest_results[runtime_item.manifest_index] = build_manifest_item(
                                runtime_item,
                                status="processed",
                                resume_status=item_resume_metadata[runtime_item.manifest_index].status,
                                resume_reason=item_resume_metadata[runtime_item.manifest_index].reason,
                                warnings=item_warnings.get(runtime_item.manifest_index),
                                processing_metadata=item_processing_metadata.get(runtime_item.manifest_index),
                            )

                    fill_prepare_queue()

    manifest_items = [item for item in manifest_results if item is not None]

    if (
        source_spec.source_kind == "path"
        and source_spec.source_root is not None
        and batch
        and failed == 0
        and processed + skipped == len(batch)
        and not disable_source_rename
    ):
        source_path_after = build_done_source_path(source_spec.source_root)
        if source_path_after != source_spec.source_root:
            try:
                mega_client.rename_path(source_spec.source_root, source_path_after)
                rename_applied = True
            except Exception as exc:  # pragma: no cover - exercised in tests via fakes.
                rename_error = str(exc)
                source_path_after = source_spec.source_root

    if source_path_before is not None and source_path_after != source_path_before:
        for item in manifest_items:
            item["source"] = rewrite_mega_path_prefix(item["source"], source_path_before, source_path_after)
            item["output_path"] = rewrite_mega_path_prefix(item["output_path"], source_path_before, source_path_after)
            item["companion_output_path"] = rewrite_mega_path_prefix(
                item["companion_output_path"],
                source_path_before,
                source_path_after,
            )

    manifest_s4_key = None
    if s4_enabled:
        manifest_s4_key = build_s4_manifest_key(
            source_spec.source_label,
            run_id,
            manifest_prefix=s4_manifest_prefix,
        )
    manifest_local_path = artifact_root / f"run-manifest-{run_id}.json"
    status = "completed" if failed == 0 and rename_error is None else "partial_failed"
    manifest = {
        "source_kind": source_spec.source_kind,
        "source": source_spec.source_value,
        "source_path_before": str(source_path_before) if source_path_before is not None else None,
        "source_path_after": str(source_path_after) if source_path_after is not None else None,
        "output_root": str(source_path_after if source_spec.source_kind == "path" else source_spec.output_root),
        "rename_applied": rename_applied,
        "rename_error": rename_error,
        "rename_deferred": bool(disable_source_rename and source_spec.source_kind == "path"),
        "resume_source": "mega",
        "resume_mode": "force" if force else "automatic",
        "resume_summary": resume_summary,
        "selected_file_kinds": sorted(selected_file_kinds),
        "compute_profile": compute_profile,
        "selection_policy": selection_policy,
        "preferred_regions": normalized_preferred_regions,
        "selected_region": selected_region,
        "instance_type": instance_type,
        "image_id": selected_image_id,
        "image_family": image_family,
        "image_version": image_version,
        "architecture": architecture,
        "gpu_count": gpu_count,
        "resolved_runtime_root": resolved_runtime_root,
        "attached_volume_mount": attached_volume_mount,
        "scheduler_batch_id": scheduler_batch_id,
        "scheduler_policy": scheduler_policy,
        "worker_name": worker_name,
        "runner_host": runner_host,
        "assigned_file_kinds": sorted(selected_file_kinds),
        "model": model,
        "s4_bucket": s4_bucket if s4_enabled else None,
        "status": status,
        "summary": {
            "processed": processed,
            "skipped": skipped,
            "failed": failed,
            "discovered": len(batch),
        },
        "warnings": manifest_warnings,
        "items": manifest_items,
    }
    write_manifest(manifest, manifest_local_path)
    if s4_enabled and manifest_s4_key is not None:
        try:
            s4_client.upload_file(manifest_local_path, s4_bucket, manifest_s4_key)
        except Exception as exc:
            manifest_warnings.append(build_s4_warning("s4 manifest upload", manifest_s4_key, exc))
            write_manifest(manifest, manifest_local_path)

    return {
        "processed": processed,
        "skipped": skipped,
        "failed": failed,
        "status": status,
        "artifact_dir": artifact_root,
        "manifest_path": manifest_local_path,
        "manifest_s4_key": manifest_s4_key,
        "manifest": manifest,
        "run_id": run_id,
        "source_label": source_spec.source_label,
        "source_path_before": source_path_before,
        "source_path_after": source_path_after,
        "rename_applied": rename_applied,
        "rename_error": rename_error,
        "rename_deferred": bool(disable_source_rename and source_spec.source_kind == "path"),
        "resume_source": "mega",
        "resume_mode": "force" if force else "automatic",
        "resume_summary": resume_summary,
        "selected_file_kinds": sorted(selected_file_kinds),
    }


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser()
    source_group = parser.add_mutually_exclusive_group(required=True)
    source_group.add_argument("--mega-source-path")
    source_group.add_argument("--mega-source-link")
    parser.add_argument("--mega-output-path")
    parser.add_argument("--s4-endpoint-url")
    parser.add_argument("--s4-bucket")
    parser.add_argument("--s4-staging-prefix", default="staging")
    parser.add_argument("--s4-transcript-prefix", default="transcripts")
    parser.add_argument("--s4-manifest-prefix", default="manifests")
    parser.add_argument("--work-dir", default=str(Path.cwd() / ".tmp-transcribe"))
    parser.add_argument("--model", default="large-v3")
    parser.add_argument("--io-workers", type=int, default=4)
    parser.add_argument("--transcribe-workers", type=int, default=1)
    parser.add_argument("--include-file-kinds")
    parser.add_argument("--disable-source-rename", action="store_true")
    parser.add_argument("--compute-profile")
    parser.add_argument("--selection-policy")
    parser.add_argument("--preferred-regions")
    parser.add_argument("--selected-region")
    parser.add_argument("--instance-type")
    parser.add_argument("--selected-image-id")
    parser.add_argument("--image-family")
    parser.add_argument("--image-version")
    parser.add_argument("--architecture")
    parser.add_argument("--gpu-count", type=int)
    parser.add_argument("--resolved-runtime-root")
    parser.add_argument("--attached-volume-mount")
    parser.add_argument("--scheduler-batch-id")
    parser.add_argument("--scheduler-policy")
    parser.add_argument("--worker-name")
    parser.add_argument("--runner-host")
    parser.add_argument("--overlap-io-with-transcription", action="store_true")
    parser.add_argument("--prefetch-root", help="NFS directory where mega_prefetch.py writes pre-downloaded files")
    parser.add_argument("--file-slice", help="Slice spec 'INDEX/TOTAL' for distributing files across GPUs, e.g. '0/16' means this worker handles files where index %% 16 == 0")
    parser.add_argument("--force", action="store_true")
    parser.add_argument("--run-id")
    return parser


def main(argv: list[str] | None = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)
    if bool(args.s4_endpoint_url) != bool(args.s4_bucket):
        parser.error("--s4-endpoint-url and --s4-bucket must be provided together.")
    s4_client = None
    if args.s4_endpoint_url:
        s4_client = S4Client.from_env(endpoint_url=args.s4_endpoint_url)
    result = run_batch(
        mega_client=MegaCli(),
        s4_client=s4_client,
        transcriber=FasterWhisperTranscriber(),
        mega_source_path=args.mega_source_path,
        mega_source_link=args.mega_source_link,
        mega_output_path=args.mega_output_path,
        work_dir=Path(args.work_dir),
        s4_bucket=args.s4_bucket,
        s4_staging_prefix=args.s4_staging_prefix,
        s4_transcript_prefix=args.s4_transcript_prefix,
        s4_manifest_prefix=args.s4_manifest_prefix,
        model=args.model,
        io_workers=args.io_workers,
        transcribe_workers=args.transcribe_workers,
        include_file_kinds=args.include_file_kinds,
        disable_source_rename=args.disable_source_rename,
        compute_profile=args.compute_profile,
        selection_policy=args.selection_policy,
        preferred_regions=args.preferred_regions,
        selected_region=args.selected_region,
        instance_type=args.instance_type,
        selected_image_id=args.selected_image_id,
        image_family=args.image_family,
        image_version=args.image_version,
        architecture=args.architecture,
        gpu_count=args.gpu_count,
        resolved_runtime_root=args.resolved_runtime_root,
        attached_volume_mount=args.attached_volume_mount,
        scheduler_batch_id=args.scheduler_batch_id,
        scheduler_policy=args.scheduler_policy,
        worker_name=args.worker_name,
        runner_host=args.runner_host,
        overlap_io_with_transcription=args.overlap_io_with_transcription,
        prefetch_root=args.prefetch_root,
        file_slice=args.file_slice,
        force=args.force,
        timestamp_value=args.run_id,
    )
    print(
        json.dumps(
            {
                "processed": result["processed"],
                "skipped": result["skipped"],
                "failed": result["failed"],
                "status": result["status"],
                "artifact_dir": str(result["artifact_dir"]),
                "manifest_path": str(result["manifest_path"]),
                "manifest_s4_key": result["manifest_s4_key"],
                "run_id": result["run_id"],
                "source_path_before": (
                    str(result.get("source_path_before")) if result.get("source_path_before") else None
                ),
                "source_path_after": (
                    str(result.get("source_path_after")) if result.get("source_path_after") else None
                ),
                "rename_applied": result.get("rename_applied"),
            },
            indent=2,
        )
    )
    return 0 if result["status"] == "completed" else 1


if __name__ == "__main__":
    raise SystemExit(main())
